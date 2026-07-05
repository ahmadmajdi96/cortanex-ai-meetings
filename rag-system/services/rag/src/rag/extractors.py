from __future__ import annotations

import csv
import io
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


@dataclass(frozen=True)
class TextSection:
    text: str
    metadata: dict[str, Any] = field(default_factory=dict)


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".md", ".markdown", ".html", ".htm", ".csv", ".xlsx", ".pptx"}


def supported_file(filename: str) -> bool:
    return Path(filename.lower()).suffix in SUPPORTED_EXTENSIONS


def _clean(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _decode(data: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "cp1256", "windows-1252", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def extract_text(filename: str, content_type: str | None, data: bytes) -> list[TextSection]:
    suffix = Path(filename.lower()).suffix
    if suffix == ".pdf" or content_type == "application/pdf":
        return _extract_pdf(data)
    if suffix == ".docx":
        return _extract_docx(data)
    if suffix in {".html", ".htm"}:
        return _extract_html(data)
    if suffix == ".csv":
        return _extract_csv(data)
    if suffix == ".xlsx":
        return _extract_xlsx(data)
    if suffix == ".pptx":
        return _extract_pptx(data)
    if suffix in {".txt", ".md", ".markdown"} or (content_type or "").startswith("text/"):
        text = _clean(_decode(data))
        return [TextSection(text=text, metadata={"section": "body"})] if text else []
    raise ValueError(f"unsupported file type for {filename!r}")


def _extract_pdf(data: bytes) -> list[TextSection]:
    reader = PdfReader(io.BytesIO(data))
    sections: list[TextSection] = []
    for page_index, page in enumerate(reader.pages, start=1):
        text = _clean(page.extract_text() or "")
        if text:
            sections.append(TextSection(text=text, metadata={"page": page_index}))
    return sections


def _extract_docx(data: bytes) -> list[TextSection]:
    document = DocxDocument(io.BytesIO(data))
    blocks: list[str] = []
    for paragraph in document.paragraphs:
        text = _clean(paragraph.text)
        if text:
            blocks.append(text)
    for table in document.tables:
        for row in table.rows:
            cells = [_clean(cell.text) for cell in row.cells if _clean(cell.text)]
            if cells:
                blocks.append(" | ".join(cells))
    text = _clean("\n\n".join(blocks))
    return [TextSection(text=text, metadata={"section": "body"})] if text else []


def _extract_html(data: bytes) -> list[TextSection]:
    soup = BeautifulSoup(_decode(data), "lxml")
    for element in soup(["script", "style", "noscript"]):
        element.decompose()
    text = _clean(soup.get_text("\n"))
    return [TextSection(text=text, metadata={"section": "body"})] if text else []


def _extract_csv(data: bytes) -> list[TextSection]:
    text = _decode(data)
    rows = []
    reader = csv.reader(io.StringIO(text))
    for row in reader:
        rows.append(" | ".join(cell.strip() for cell in row if cell.strip()))
    clean = _clean("\n".join(row for row in rows if row))
    return [TextSection(text=clean, metadata={"section": "csv"})] if clean else []


def _extract_xlsx(data: bytes) -> list[TextSection]:
    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sections: list[TextSection] = []
    for sheet in workbook.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if values:
                rows.append(" | ".join(values))
        text = _clean("\n".join(rows))
        if text:
            sections.append(TextSection(text=text, metadata={"sheet": sheet.title}))
    return sections


def _extract_pptx(data: bytes) -> list[TextSection]:
    presentation = Presentation(io.BytesIO(data))
    sections: list[TextSection] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = _clean(shape.text)
                if text:
                    parts.append(text)
        text = _clean("\n".join(parts))
        if text:
            sections.append(TextSection(text=text, metadata={"slide": slide_index}))
    return sections
